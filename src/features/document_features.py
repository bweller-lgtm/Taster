"""Document feature extraction and similarity grouping.

Supports PDF, DOCX, XLSX, PPTX, HTML, TXT, MD, CSV, and RTF formats.
"""
import os
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, field

import numpy as np


@dataclass
class DocumentFeatures:
    """Extracted features from a document."""
    text_content: str = ""                    # Extracted text (truncated)
    text_length: int = 0                      # Total character count
    page_count: Optional[int] = None          # For PDFs, PPTXs
    word_count: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)
    file_format: str = ""                     # pdf, docx, xlsx, etc.
    has_images: bool = False
    embedding: Optional[np.ndarray] = None    # Text embedding for similarity
    file_size_bytes: int = 0


class DocumentFeatureExtractor:
    """Extract features from supported document formats."""

    def __init__(self, max_chars: int = 50000, max_pages: int = 20):
        self.max_chars = max_chars
        self.max_pages = max_pages
        self._embedding_model = None

    def extract_features(self, doc_path: Path) -> DocumentFeatures:
        """Extract features from any supported document."""
        doc_path = Path(doc_path)
        ext = doc_path.suffix.lower()

        features = DocumentFeatures(
            file_format=ext.lstrip("."),
            file_size_bytes=doc_path.stat().st_size,
        )

        # Extract text based on format
        try:
            text = self.extract_text(doc_path)
            features.text_content = text[:self.max_chars]
            features.text_length = len(text)
            features.word_count = len(text.split())
        except Exception as e:
            print(f"Warning: Text extraction failed for {doc_path.name}: {e}")

        # Extract metadata
        try:
            features.metadata = self.extract_metadata(doc_path)
            features.page_count = features.metadata.get("page_count")
            features.has_images = features.metadata.get("has_images", False)
        except Exception as e:
            print(f"Warning: Metadata extraction failed for {doc_path.name}: {e}")

        return features

    def extract_text(self, doc_path: Path) -> str:
        """Extract text content from document, dispatching by format."""
        ext = doc_path.suffix.lower()
        extractors = {
            ".pdf": self._extract_pdf_text,
            ".docx": self._extract_docx_text,
            ".xlsx": self._extract_xlsx_text,
            ".pptx": self._extract_pptx_text,
            ".html": self._extract_html_text,
            ".htm": self._extract_html_text,
            ".txt": self._extract_plain_text,
            ".md": self._extract_plain_text,
            ".csv": self._extract_plain_text,
            ".rtf": self._extract_plain_text,
        }
        extractor = extractors.get(ext, self._extract_plain_text)
        return extractor(doc_path)

    def extract_metadata(self, doc_path: Path) -> Dict[str, Any]:
        """Extract metadata from document."""
        ext = doc_path.suffix.lower()
        meta = {
            "filename": doc_path.name,
            "extension": ext,
            "file_size_bytes": doc_path.stat().st_size,
        }

        try:
            if ext == ".pdf":
                meta.update(self._extract_pdf_metadata(doc_path))
            elif ext == ".docx":
                meta.update(self._extract_docx_metadata(doc_path))
            elif ext == ".xlsx":
                meta.update(self._extract_xlsx_metadata(doc_path))
            elif ext == ".pptx":
                meta.update(self._extract_pptx_metadata(doc_path))
        except Exception:
            pass

        return meta

    def compute_text_embedding(self, text: str) -> Optional[np.ndarray]:
        """Compute sentence-transformer embedding for similarity."""
        if not text.strip():
            return None

        try:
            model = self._get_embedding_model()
            if model is None:
                return None
            # Truncate to reasonable length for embedding
            truncated = text[:10000]
            embedding = model.encode(truncated, normalize_embeddings=True)
            return np.array(embedding)
        except Exception as e:
            print(f"Warning: Embedding computation failed: {e}")
            return None

    def compute_similarity(self, features_a: DocumentFeatures, features_b: DocumentFeatures) -> float:
        """Compute text similarity between two documents."""
        if features_a.embedding is None or features_b.embedding is None:
            return 0.0
        return float(np.dot(features_a.embedding, features_b.embedding))

    def _get_embedding_model(self):
        """Lazily load the sentence-transformers model."""
        if self._embedding_model is None:
            try:
                from sentence_transformers import SentenceTransformer
                self._embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
            except ImportError:
                print("Warning: sentence-transformers not installed. Text embeddings disabled.")
                return None
        return self._embedding_model

    # --- Format-specific text extractors ---

    def _extract_pdf_text(self, path: Path) -> str:
        """Extract text from PDF."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = reader.pages[:self.max_pages]
            texts = []
            for page in pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
            return "\n\n".join(texts)
        except ImportError:
            return self._extract_plain_text(path)
        except Exception as e:
            print(f"Warning: PDF extraction failed for {path.name}: {e}")
            return ""

    def _extract_docx_text(self, path: Path) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            return ""
        except Exception as e:
            print(f"Warning: DOCX extraction failed for {path.name}: {e}")
            return ""

    def _extract_xlsx_text(self, path: Path) -> str:
        """Extract text from Excel spreadsheet."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames[:10]:
                sheet = wb[sheet_name]
                texts.append(f"--- Sheet: {sheet_name} ---")
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if row_count >= 100:
                        texts.append("... (truncated)")
                        break
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        texts.append("\t".join(cells))
                    row_count += 1
            wb.close()
            return "\n".join(texts)
        except ImportError:
            return ""
        except Exception as e:
            print(f"Warning: XLSX extraction failed for {path.name}: {e}")
            return ""

    def _extract_pptx_text(self, path: Path) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                if slide_text:
                    texts.append(f"--- Slide {i+1} ---")
                    texts.extend(slide_text)
            return "\n".join(texts)
        except ImportError:
            return ""
        except Exception as e:
            print(f"Warning: PPTX extraction failed for {path.name}: {e}")
            return ""

    def _extract_html_text(self, path: Path) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            soup = BeautifulSoup(content, "html.parser")
            # Remove script and style elements
            for element in soup(["script", "style"]):
                element.decompose()
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            # Fall back to plain text
            return self._extract_plain_text(path)
        except Exception as e:
            print(f"Warning: HTML extraction failed for {path.name}: {e}")
            return ""

    def _extract_plain_text(self, path: Path) -> str:
        """Extract text from plain text file."""
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        for encoding in encodings:
            try:
                with open(path, "r", encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        return ""

    # --- Format-specific metadata extractors ---

    def _extract_pdf_metadata(self, path: Path) -> Dict[str, Any]:
        """Extract metadata from PDF."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            meta = {"page_count": len(reader.pages)}
            if reader.metadata:
                if reader.metadata.title:
                    meta["title"] = reader.metadata.title
                if reader.metadata.author:
                    meta["author"] = reader.metadata.author
                if reader.metadata.creation_date:
                    meta["creation_date"] = str(reader.metadata.creation_date)
            return meta
        except ImportError:
            return {}

    def _extract_docx_metadata(self, path: Path) -> Dict[str, Any]:
        """Extract metadata from Word document."""
        try:
            from docx import Document
            doc = Document(str(path))
            meta = {}
            props = doc.core_properties
            if props.title:
                meta["title"] = props.title
            if props.author:
                meta["author"] = props.author
            if props.created:
                meta["creation_date"] = str(props.created)
            meta["paragraph_count"] = len(doc.paragraphs)
            meta["has_images"] = any(
                r.element.tag.endswith("}blipFill")
                for p in doc.paragraphs
                for r in p.runs
                if r.element is not None
                for child in r.element
            )
            return meta
        except ImportError:
            return {}
        except Exception:
            return {}

    def _extract_xlsx_metadata(self, path: Path) -> Dict[str, Any]:
        """Extract metadata from Excel spreadsheet."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True)
            meta = {
                "sheet_names": wb.sheetnames,
                "sheet_count": len(wb.sheetnames),
            }
            wb.close()
            return meta
        except ImportError:
            return {}

    def _extract_pptx_metadata(self, path: Path) -> Dict[str, Any]:
        """Extract metadata from PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            meta = {"page_count": len(prs.slides)}
            props = prs.core_properties
            if props.title:
                meta["title"] = props.title
            if props.author:
                meta["author"] = props.author
            has_images = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True
                        break
                if has_images:
                    break
            meta["has_images"] = has_images
            return meta
        except ImportError:
            return {}


class DocumentGrouper:
    """Groups similar documents using text embeddings."""

    def __init__(self, similarity_threshold: float = 0.85):
        self.similarity_threshold = similarity_threshold

    def group_documents(
        self,
        documents: List[Path],
        features: Dict[Path, DocumentFeatures],
    ) -> List[List[Path]]:
        """Group similar documents for comparative classification.

        Returns list of groups (singletons for unique docs).
        """
        if not documents:
            return []

        # Build similarity matrix using embeddings
        embeddings = []
        valid_docs = []
        for doc in documents:
            feat = features.get(doc)
            if feat and feat.embedding is not None:
                embeddings.append(feat.embedding)
                valid_docs.append(doc)
            else:
                # No embedding available - treat as singleton
                pass

        if not embeddings:
            return [[doc] for doc in documents]

        embeddings_matrix = np.stack(embeddings)

        # Compute pairwise similarity
        similarity = np.dot(embeddings_matrix, embeddings_matrix.T)

        # Union-find grouping
        n = len(valid_docs)
        parent = list(range(n))

        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x

        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb

        for i in range(n):
            for j in range(i + 1, n):
                if similarity[i, j] >= self.similarity_threshold:
                    union(i, j)

        # Build groups
        groups_map: Dict[int, List[Path]] = {}
        for i, doc in enumerate(valid_docs):
            root = find(i)
            if root not in groups_map:
                groups_map[root] = []
            groups_map[root].append(doc)

        groups = list(groups_map.values())

        # Add docs without embeddings as singletons
        docs_with_embeddings = set(valid_docs)
        for doc in documents:
            if doc not in docs_with_embeddings:
                groups.append([doc])

        return groups
